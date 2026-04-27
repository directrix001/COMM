"""
tabs/tab7_ppt_upload.py
────────────────────────
Tab 7 — PPT Upload & Push to Master Database

Mirrors the PyQt6 Upload tab logic in Streamlit:
  1. User uploads a PowerPoint (.pptx / .ppt)
  2. Text is extracted slide-by-slide, headers are parsed, and
     metadata columns (Category, Scenarios, Functions, Month, Year …)
     are derived via regex.
  3. Five keyword-extraction passes enrich each row:
       CostCat description, Function_desc, Entity_desc, Criteria, Region
  4. The enriched DataFrame is shown in an editable data editor.
  5. "Push to Master Database" appends + deduplicates into
     database/Segregateddata.xlsx
  6. "Download File" saves the current editor state as Excel.

Public API
----------
render()   — call inside  `with tab:` in app.py
"""

from __future__ import annotations

import io
import os
import re
from datetime import datetime
from pathlib import Path

import pandas as pd
import streamlit as st

# ── Path constant (same as tab6) ─────────────────────────────────────────────
MASTER_XLSX_PATH = Path("database") / "Segregateddata.xlsx"


# ══════════════════════════════════════════════════════════════════════════════
# REFERENCE LISTS  (ported 1-to-1 from the PyQt6 source)
# ══════════════════════════════════════════════════════════════════════════════

REGIONS_AND_COUNTRIES = [
    'Afghanistan','Albania','Algeria','Andorra','Angola',
    'Antigua and Barbuda','Argentina','Armenia','Australia','Austria',
    'Azerbaijan','Bahamas','Bahrain','Bangladesh','Barbados',
    'Belarus','Belgium','Belize','Benin','Bhutan',
    'Bolivia','Bosnia and Herzegovina','Botswana','Brazil','Brunei',
    'Bulgaria','Burkina Faso','Burundi','Cabo Verde','Cambodia',
    'Cameroon','Canada','Central African Republic','Chad','Chile',
    'China','Colombia','Comoros','Congo (Congo-Brazzaville)','Costa Rica',
    'Croatia','Cuba','Cyprus','Czechia (Czech Republic)','Democratic Republic of the Congo',
    'Denmark','Djibouti','Dominica','Dominican Republic','Ecuador',
    'Egypt','El Salvador','Equatorial Guinea','Eritrea','Estonia',
    'Eswatini (fmr. "Swaziland")','Ethiopia','Fiji','Finland','France',
    'Gabon','Gambia','Georgia','Germany','Ghana',
    'Greece','Grenada','Guatemala','Guinea','Guinea-Bissau',
    'Guyana','Haiti','Holy See','Honduras','Hungary',
    'Iceland','India','Indonesia','Iran','Iraq',
    'Ireland','Israel','Italy','Jamaica','Japan',
    'Jordan','Kazakhstan','Kenya','Kiribati','Kuwait',
    'Kyrgyzstan','Laos','Latvia','Lebanon','Lesotho',
    'Liberia','Libya','Liechtenstein','Lithuania','Luxembourg',
    'Madagascar','Malawi','Malaysia','Maldives','Mali',
    'Malta','Marshall Islands','Mauritania','Mauritius','Mexico',
    'Micronesia','Moldova','Monaco','Mongolia','Montenegro',
    'Morocco','Mozambique','Myanmar','Namibia','Nauru',
    'Nepal','Netherlands','New Zealand','Nicaragua','Niger',
    'Nigeria','North Korea','North Macedonia','Norway','Oman',
    'Pakistan','Palau','Palestine State','Panama','Papua New Guinea',
    'Paraguay','Peru','Philippines','Poland','Portugal',
    'Qatar','Romania','Russia','Rwanda','Saint Kitts and Nevis',
    'Saint Lucia','Saint Vincent and the Grenadines','Samoa','San Marino','Sao Tome and Principe',
    'Saudi Arabia','Senegal','Serbia','Seychelles','Sierra Leone',
    'Singapore','Slovakia','Slovenia','Solomon Islands','Somalia',
    'South Africa','South Korea','South Sudan','Spain','Sri Lanka',
    'Sudan','Suriname','Sweden','Switzerland','Syria',
    'Tajikistan','Tanzania','Thailand','Timor-Leste','Togo',
    'Tonga','Trinidad and Tobago','Tunisia','Turkey','Turkmenistan',
    'Tuvalu','Uganda','Ukraine','United Arab Emirates','United Kingdom',
    'United States of America','Uruguay','Uzbekistan','Vanuatu','Venezuela',
    'Vietnam','Yemen','Zambia','Zimbabwe',
    'AMIO','AMIEO','Middle East','Europe','Oceania','ME','Africa',
]

LC_OH = ['LC', 'OH', 'Notes', 'OVH']

KEYWORD_1 = [
    "A&F SVP OFFICE","Milky Way","CUSTOMS","Finance Project portfolio",
    "FINANCE & ACCOUNTING","FPLAN & VP FINANCE","NRBS FINANCE","NRBS MD Office",
    "NRBS TAX","PEC","S&M CONTROLLING FIN","SSC","TAX","TREASURY",
    "R&D Finance & Accounting","IS","NRBSISIT","CENTRAL TASKS","CEC",
    "GLOBAL ROLES","QEX","NON EUROPE & EXEC","CPLO SVP Office","General Affairs",
    "NRBSGA","Regional General Affairs","Local HR","Regional COE & HRBP","NRBS HR",
    "Local Health & Safety","Regional Health & Safety","REGIONAL SECURITY",
    "LOCAL SECURITY","Business Transformation and Acceleration","Corporate projects",
    "STRATEGIC PLANNING","INT AUDIT & JSOX","MC-E CHAIRMAN OFFICE","ELECTRIFICATION",
    "OPD","MZK SVP Office","Central SCM","NRBS Purchasing","PURCHASING",
    "AS Import & Service repairs","AFTER SALES OPS","REGION - AFTER SALES",
    "Service Marketing","Brand & Customer Experience","Customer Exp 2.0","MARCOM",
    "CONNECTED SERVICES","MARKET INTELLIGENCE","RBU - MARKETING",
    "Business Development Europe","DX for CX","FULL LIFE CYCLE",
    "REGION - Energy Services","Transformation","PRODUCT STRATEGY & PLANNING",
    "RBU - DND & CQ","RBU - SALES","RBU - AFTER SALES","RBU - ENERGY AND BATTERY",
    "RBU MD office","Product & Service Planning","REGION - DND","DVP OFFICE",
    "REGION - INSC","MARKETING OPS","SALES OPS","LCV","NANO",
    "Planning & Performance","M&S - SVP OFFICE","COMMUNICATION","COMPLIANCE",
    "EGA (External & Governmental Affairs)","LEGAL","SCAG_SVP Office",
    "Sustainability & Governance","TCS","TCS AVES",
]

CRITERIA = [
    "PROCURED SERVICES","TRAVEL & MEALS","EMPLOYEE WELFARE",
    "RECHARGE NISSAN Level0","OPERATING COSTS","OFFICE SPACE",
    "EMPLOYEE ACTIVITY COSTS","TAX","RECHARGE OUTSIDE",
    "PROVISION FOR DOUBTFUL DEBTS","COMPANY CAR COSTS","DEPRECIATION",
]

ENTITY = [
    "Nissan Automotive Europe","NIBSA","NMISA","NMUK - PLANT",
    "Nissan Motor Parts Center","NWE - Belgium branch","NCE Germany",
    "NCE - Swiss branch","NCE - Austria branch","NITA",
    "NWE - Netherlands branch","NMGB","NNE","NNE - Denmark branch",
    "NNE - Sweden branch","NNE - Norway branch","Nissan Automobiles Turkey (NOAS)",
    "NM Ukraine Nissan/Datsun","NSCEE - Hungarian branch","NSCEE - Czech branch",
    "NSCEE - Polish branch","NMUK - NTCE","Nissan France","Nissan International SA",
    "NTCE Germany Branch","NTCE Belgium Branch","Adjustment entity",
    "NMPC - UK Lutterworth","NMUK","NAE","NRBS","NTCE","NRBS","NMEF",
]


# ══════════════════════════════════════════════════════════════════════════════
# EXTRACTION HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def _match_list(text: str, lst: list[str]) -> str:
    if not isinstance(text, str) or not text:
        return ""
    tl = text.lower()
    hits = [item for item in lst if re.search(r'\b' + re.escape(item.lower()) + r'\b', tl)]
    return ", ".join(hits)


def _extract_region(text: str) -> str:
    return _match_list(text, REGIONS_AND_COUNTRIES)

def _extract_lc_oh(text: str) -> str:
    return _match_list(text, LC_OH)

def _extract_criteria(text: str) -> str:
    return _match_list(text, CRITERIA)

def _extract_keyword(text: str) -> str:
    return _match_list(text, KEYWORD_1)

def _extract_entity(text: str) -> str:
    return _match_list(text, ENTITY)


# ══════════════════════════════════════════════════════════════════════════════
# PPT PROCESSING  (ported from process_ppt_and_spreadsheet)
# ══════════════════════════════════════════════════════════════════════════════

def _process_ppt(file_obj, filename: str) -> pd.DataFrame:
    """
    Extract text from every slide, derive metadata from the header line
    that starts with 'G&A Evolution', then run keyword-extraction passes.
    Returns the enriched DataFrame (rows with Comments longer than 35 chars).
    """
    from pptx import Presentation

    prs = Presentation(file_obj)
    extracted_content: list[dict] = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        header = None

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()

                if (
                    header is None
                    and text.startswith("G&A Evolution")
                    and text not in ["AMIEO", "M€ @BP25FX"]
                ):
                    header = text

                extracted_content.append(
                    {"Slide Number": slide_number, "Comments": text}
                )

        if header:
            header = header.replace("-", "–")
            header_parts = header.split(" – ")

            category       = header_parts[0] if len(header_parts) > 0 else ""
            month          = ""
            scenario       = ""
            function_      = ""
            function_view  = ""

            month_match = re.search(
                r"\b(January|February|March|April|May|June|July|August|September"
                r"|October|November|December|Jan|Feb|Mar|Apr|May|Jun|Jul|Aug"
                r"|Sep|Oct|Nov|Dec)\b",
                header, re.IGNORECASE,
            )
            if month_match:
                month = month_match.group(0)

            scenario_match = re.search(r"(MTD|YTD)\s+vs\.\s+[^–]+", header)
            if scenario_match:
                scenario = scenario_match.group(0)

            if len(header_parts) > 2:
                function_ = header_parts[2]
            if len(header_parts) > 3:
                function_view = header_parts[3]

            current_year = datetime.now().year

            for item in extracted_content:
                if item["Slide Number"] == slide_number:
                    item.update(
                        {
                            "File_name":      filename,
                            "Category":       category,
                            "Scenarios":      scenario,
                            "Functions":      function_,
                            "Functions-View": function_view,
                            "Month":          month,
                            "Year":           current_year,
                            "Forecast":       "",
                            "Actual":         "",
                            "Variance":       "",
                        }
                    )

    # Keep only rows with meaningful comment length
    filtered = [i for i in extracted_content if len(i.get("Comments", "")) > 35]
    df = pd.DataFrame(filtered)

    if df.empty:
        return df

    # Keyword-extraction columns
    df["CostCat description"] = df["Comments"].apply(_extract_criteria)
    df["Function_desc"]       = df["Comments"].apply(_extract_keyword)
    df["Entity_desc"]         = df["Comments"].apply(_extract_entity)
    df["Criteria"]            = df["Comments"].apply(_extract_lc_oh)
    df["Region"]              = df["Comments"].apply(_extract_region)

    return df


# ══════════════════════════════════════════════════════════════════════════════
# PUSH HELPER
# ══════════════════════════════════════════════════════════════════════════════

def _push_to_master(df_new: pd.DataFrame) -> str:
    """Append df_new to master, dedup, save.  Returns status message."""
    MASTER_XLSX_PATH.parent.mkdir(parents=True, exist_ok=True)

    if MASTER_XLSX_PATH.exists():
        master_df = pd.read_excel(MASTER_XLSX_PATH)
    else:
        master_df = pd.DataFrame()

    combined = (
        pd.concat([master_df, df_new], ignore_index=True)
        if not master_df.empty
        else df_new.copy()
    )

    if "Slide Number" in combined.columns and "Comments" in combined.columns:
        combined = combined.drop_duplicates(
            subset=["Slide Number", "Comments"], keep="last"
        ).reset_index(drop=True)

    combined.to_excel(MASTER_XLSX_PATH, index=False)
    return f"✅ Pushed {len(df_new):,} rows → master now has {len(combined):,} rows."


# ══════════════════════════════════════════════════════════════════════════════
# MAIN RENDER
# ══════════════════════════════════════════════════════════════════════════════

def render() -> None:
    st.markdown(
        '<p class="va-section-label">📂 PPT Upload & Master Database Push</p>',
        unsafe_allow_html=True,
    )
    st.markdown(
        "Upload the **current month SVP PowerPoint deck**. The tool extracts "
        "slide comments, derives metadata, lets you review & edit the rows, "
        "then pushes them to the master commentary database."
    )

    st.markdown('<hr class="va-divider">', unsafe_allow_html=True)

    # ── File uploader ────────────────────────────────────────────────────────
    uploaded = st.file_uploader(
        "Upload PowerPoint file (.pptx / .ppt)",
        type=["pptx", "ppt"],
        key="ppt_upload_file",
    )

    if uploaded is None:
        st.info("Upload a PowerPoint file above to begin extraction.")
        return

    # ── Process (cached per file bytes) ──────────────────────────────────────
    file_bytes = uploaded.read()

    @st.cache_data(show_spinner="Extracting text from slides…")
    def _cached_process(b: bytes, fname: str) -> pd.DataFrame:
        return _process_ppt(io.BytesIO(b), fname)

    try:
        extracted_df = _cached_process(file_bytes, uploaded.name)
    except Exception as exc:
        st.error(f"PPT extraction failed: {exc}")
        return

    if extracted_df.empty:
        st.warning("No comments longer than 35 characters were found in this file.")
        return

    st.success(
        f"✅ Extracted **{len(extracted_df):,} rows** from "
        f"**{uploaded.name}** — review & edit below before pushing."
    )

    st.markdown('<hr class="va-divider">', unsafe_allow_html=True)
    st.markdown(
        '<p class="va-section-label">✏️ Review & Edit Extracted Data</p>',
        unsafe_allow_html=True,
    )

    # ── Editable data editor ─────────────────────────────────────────────────
    edited_df = st.data_editor(
        extracted_df,
        use_container_width=True,
        num_rows="dynamic",
        height=420,
        key="ppt_editor",
    )

    st.markdown('<hr class="va-divider">', unsafe_allow_html=True)

    btn_col1, btn_col2 = st.columns(2)

    # ── Push to master ────────────────────────────────────────────────────────
    with btn_col1:
        if st.button("🚀 Push to Master Database", key="ppt_push_btn"):
            try:
                msg = _push_to_master(edited_df)
                st.success(msg)
                # Bust the search-tab cache by touching the file mtime
                st.cache_data.clear()
                st.toast("Master database updated!", icon="✅")
            except Exception as exc:
                st.error(f"Push failed: {exc}")

    # ── Download edited data ──────────────────────────────────────────────────
    with btn_col2:
        xl_buf = io.BytesIO()
        with pd.ExcelWriter(xl_buf, engine="openpyxl") as w:
            edited_df.to_excel(w, index=False, sheet_name="Commentary_edited")
        st.download_button(
            "📥 Download Edited File (.xlsx)",
            data=xl_buf.getvalue(),
            file_name="Commentary_edited.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            key="ppt_dl_btn",
        )

    # ── Preview of current master DB ─────────────────────────────────────────
    st.markdown('<hr class="va-divider">', unsafe_allow_html=True)
    with st.expander("👁️ Preview current Master Database", expanded=False):
        if MASTER_XLSX_PATH.exists():
            master_preview = pd.read_excel(MASTER_XLSX_PATH)
            st.caption(f"{len(master_preview):,} rows in master")
            st.dataframe(master_preview.head(100), use_container_width=True, height=280)
        else:
            st.info("Master database does not exist yet. Push data to create it.")
