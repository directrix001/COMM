"""
routers/tab7_ppt.py  — FIXED
Fixes:
  1. 'Slide Number' column (space in name) crashes itertuples → use iloc
  2. NaN / inf in Year/numeric cols → not JSON compliant
  3. df.where().to_dict() misses numpy NaN → replaced with _safe_records
"""
from __future__ import annotations

import io
import math
import re
import uuid
from datetime import datetime
from pathlib import Path

import pandas as pd
from fastapi import APIRouter, File, HTTPException, Request, UploadFile
from fastapi.responses import JSONResponse, StreamingResponse

from services import session_store as ss

router = APIRouter()

MASTER_XLSX = Path("database") / "Segregateddata.xlsx"


def _sid(request: Request) -> str:
    return request.cookies.get("va_sid", str(uuid.uuid4()))


# ═══════════════════════════════════════════════════════════════
# CORE FIX — safe serialisation (no itertuples, handles spaces)
# ═══════════════════════════════════════════════════════════════

def _safe_val(v):
    """Convert NaN / inf / numpy scalars to JSON-safe Python types."""
    if v is None:
        return None
    if isinstance(v, float):
        if math.isnan(v) or math.isinf(v):
            return None
        return v
    try:
        import numpy as np
        if isinstance(v, (np.floating,)):
            f = float(v)
            return None if (math.isnan(f) or math.isinf(f)) else f
        if isinstance(v, (np.integer,)):
            return int(v)
        if isinstance(v, np.bool_):
            return bool(v)
        if isinstance(v, float) and (math.isnan(v) or math.isinf(v)):
            return None
    except ImportError:
        pass
    return v


def _safe_records(df: pd.DataFrame) -> list:
    """
    Safe DataFrame → list[dict] using iloc (not itertuples).
    itertuples() mangles column names that contain spaces or special chars
    (e.g. 'Slide Number' → '_1'), so getattr fails.
    iloc avoids that entirely.
    """
    cols = df.columns.tolist()
    result = []
    for i in range(len(df)):
        row_dict = {}
        for j, col in enumerate(cols):
            row_dict[col] = _safe_val(df.iloc[i, j])
        result.append(row_dict)
    return result


def _sanitise_df(df: pd.DataFrame) -> pd.DataFrame:
    """Replace NaN/inf in all columns; force strings where appropriate."""
    df = df.copy()

    # String columns — replace NaN with ""
    for col in df.select_dtypes(include=["object"]).columns:
        df[col] = df[col].apply(
            lambda x: "" if (
                x is None or
                (isinstance(x, float) and (math.isnan(x) or math.isinf(x)))
            ) else str(x) if not isinstance(x, str) else x
        )

    # Float columns — replace NaN/inf with None
    for col in df.select_dtypes(include=["float", "float64", "float32"]).columns:
        df[col] = df[col].apply(
            lambda x: None if (
                x is None or
                (isinstance(x, float) and (math.isnan(x) or math.isinf(x)))
            ) else x
        )

    return df


# ═══════════════════════════════════════════════════════════════
# REFERENCE LISTS
# ═══════════════════════════════════════════════════════════════

REGIONS = [
    'Afghanistan','Albania','Algeria','Andorra','Angola','Argentina','Armenia',
    'Australia','Austria','Azerbaijan','Bahrain','Bangladesh','Belgium','Brazil',
    'Bulgaria','Canada','Chile','China','Colombia','Croatia','Cuba','Cyprus',
    'Denmark','Egypt','Estonia','Ethiopia','Finland','France','Germany','Ghana',
    'Greece','Hungary','Iceland','India','Indonesia','Iran','Iraq','Ireland',
    'Israel','Italy','Jamaica','Japan','Jordan','Kazakhstan','Kenya','Kuwait',
    'Latvia','Lebanon','Lithuania','Luxembourg','Malaysia','Malta','Mexico',
    'Morocco','Netherlands','New Zealand','Nigeria','Norway','Oman','Pakistan',
    'Philippines','Poland','Portugal','Qatar','Romania','Russia','Saudi Arabia',
    'Serbia','Singapore','Slovakia','Slovenia','South Africa','South Korea',
    'Spain','Sri Lanka','Sweden','Switzerland','Syria','Thailand','Tunisia',
    'Turkey','Ukraine','United Arab Emirates','United Kingdom',
    'United States of America','Vietnam','Zimbabwe',
    'AMIO','AMIEO','Middle East','Europe','Oceania','ME','Africa',
]
LC_OH    = ['LC','OH','Notes','OVH']
CRITERIA = [
    'PROCURED SERVICES','TRAVEL & MEALS','EMPLOYEE WELFARE',
    'RECHARGE NISSAN Level0','OPERATING COSTS','OFFICE SPACE',
    'EMPLOYEE ACTIVITY COSTS','TAX','RECHARGE OUTSIDE',
    'PROVISION FOR DOUBTFUL DEBTS','COMPANY CAR COSTS','DEPRECIATION',
]
KEYWORDS = [
    'FINANCE & ACCOUNTING','TREASURY','TAX','SSC','LEGAL','COMPLIANCE',
    'COMMUNICATION','PURCHASING','AFTER SALES OPS','MARCOM',
    'CONNECTED SERVICES','MARKET INTELLIGENCE','ELECTRIFICATION','General Affairs',
]
ENTITY = [
    'Nissan Automotive Europe','NIBSA','NMISA','NMUK - PLANT',
    'NCE Germany','NITA','NMGB','NNE','Nissan France',
    'Nissan International SA','NAE','NRBS','NTCE','NMEF',
]


def _match(text: str, lst: list) -> str:
    if not isinstance(text, str) or not text:
        return ""
    tl = text.lower()
    return ", ".join(i for i in lst if re.search(r'\b' + re.escape(i.lower()) + r'\b', tl))


# ═══════════════════════════════════════════════════════════════
# PPT EXTRACTION
# ═══════════════════════════════════════════════════════════════

def _extract_ppt(file_obj, filename: str) -> pd.DataFrame:
    from pptx import Presentation

    prs     = Presentation(file_obj)
    content = []

    for slide_num, slide in enumerate(prs.slides, 1):
        header = None

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if (header is None
                        and text.startswith("G&A Evolution")
                        and text not in ["AMIEO", "M€ @BP25FX"]):
                    header = text
                content.append({"Slide_Number": slide_num, "Comments": text})
                # ↑ NOTE: renamed to Slide_Number (underscore) — no space = no itertuples crash

        if header:
            header   = header.replace("-", "–")
            parts    = header.split(" – ")
            category = parts[0] if parts else ""
            month = scenario = fn = fn_view = ""

            m = re.search(
                r"\b(January|February|March|April|May|June|July|August|September"
                r"|October|November|December|Jan|Feb|Mar|Apr|May|Jun|Jul|Aug"
                r"|Sep|Oct|Nov|Dec)\b",
                header, re.IGNORECASE,
            )
            if m:
                month = m.group(0)

            m2 = re.search(r"(MTD|YTD)\s+vs\.\s+[^–]+", header)
            if m2:
                scenario = m2.group(0)

            if len(parts) > 2:
                fn = parts[2]
            if len(parts) > 3:
                fn_view = parts[3]

            for item in content:
                if item["Slide_Number"] == slide_num:
                    item.update({
                        "File_name":      filename,
                        "Category":       category,
                        "Scenarios":      scenario,
                        "Functions":      fn,
                        "Functions-View": fn_view,
                        "Month":          month,
                        "Year":           str(datetime.now().year),  # string avoids int64
                        "Forecast":       "",
                        "Actual":         "",
                        "Variance":       "",
                    })

    # Keep only rows with meaningful comment length
    filtered = [i for i in content if len(i.get("Comments", "")) > 35]
    if not filtered:
        return pd.DataFrame()

    df = pd.DataFrame(filtered)

    # Keyword enrichment
    df["CostCat description"] = df["Comments"].apply(lambda t: _match(t, CRITERIA))
    df["Function_desc"]       = df["Comments"].apply(lambda t: _match(t, KEYWORDS))
    df["Entity_desc"]         = df["Comments"].apply(lambda t: _match(t, ENTITY))
    df["Criteria"]            = df["Comments"].apply(lambda t: _match(t, LC_OH))
    df["Region"]              = df["Comments"].apply(lambda t: _match(t, REGIONS))

    # Sanitise all columns before returning
    return _sanitise_df(df)


# ═══════════════════════════════════════════════════════════════
# ENDPOINTS
# ═══════════════════════════════════════════════════════════════

@router.post("/upload")
async def upload_ppt(request: Request, file: UploadFile = File(...)):
    sid = _sid(request)

    if not file.filename.lower().endswith((".pptx", ".ppt")):
        raise HTTPException(400, "Only .pptx / .ppt files are accepted.")

    raw = await file.read()

    try:
        df = _extract_ppt(io.BytesIO(raw), file.filename)
    except Exception as exc:
        raise HTTPException(500, f"PPT extraction failed: {exc}")

    if df.empty:
        raise HTTPException(422, "No comments longer than 35 characters found in this file.")

    # Store in session — orient=split handles None cleanly
    ss.set(sid, "ppt_extract_bytes", df.to_json(orient="split"))

    resp = JSONResponse({
        "session_id": sid,
        "rows":       len(df),
        "columns":    df.columns.tolist(),
        "records":    _safe_records(df),   # ← iloc-based, no space-column crash
    })
    resp.set_cookie("va_sid", sid, httponly=True, samesite="lax")
    return resp


@router.post("/push")
async def push_to_master(request: Request):
    sid = _sid(request)
    raw = ss.get(sid, "ppt_extract_bytes")
    if not raw:
        raise HTTPException(404, "No extracted data in session. Upload a PPT first.")

    df_new = pd.read_json(io.StringIO(raw), orient="split")
    MASTER_XLSX.parent.mkdir(parents=True, exist_ok=True)

    if MASTER_XLSX.exists():
        master = pd.read_excel(MASTER_XLSX)
    else:
        master = pd.DataFrame()

    combined = (
        pd.concat([master, df_new], ignore_index=True)
        if not master.empty
        else df_new.copy()
    )

    # Deduplicate — handle both old "Slide Number" and new "Slide_Number"
    dedup_col = "Slide_Number" if "Slide_Number" in combined.columns else "Slide Number"
    if dedup_col in combined.columns and "Comments" in combined.columns:
        combined = combined.drop_duplicates(
            subset=[dedup_col, "Comments"], keep="last"
        ).reset_index(drop=True)

    combined.to_excel(MASTER_XLSX, index=False)

    return JSONResponse({
        "status":      "ok",
        "pushed_rows": len(df_new),
        "master_rows": len(combined),
    })


@router.get("/download")
async def download(request: Request):
    sid = _sid(request)
    raw = ss.get(sid, "ppt_extract_bytes")
    if not raw:
        raise HTTPException(404, "No extracted data in session.")

    df  = pd.read_json(io.StringIO(raw), orient="split")
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as w:
        df.to_excel(w, index=False, sheet_name="Commentary_edited")
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=Commentary_edited.xlsx"},
    )


@router.get("/master")
async def master_preview():
    if not MASTER_XLSX.exists():
        return JSONResponse({"columns": [], "records": [], "total": 0})

    df      = pd.read_excel(MASTER_XLSX)
    df      = _sanitise_df(df)
    preview = df.head(100)

    return JSONResponse({
        "columns": df.columns.tolist(),
        "records": _safe_records(preview),   # ← iloc-based
        "total":   len(df),
    })
