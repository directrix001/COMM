"""
utils/ai_engine.py
──────────────────
LangGraph agent state, nodes, graph builder,
PPTX generator, and Azure env helper.
Used by Tab 3 (Commentary Generator).
"""

from __future__ import annotations

import io
import os
import re
from datetime import datetime
from typing import Any, Dict, List, Tuple

import pandas as pd
import streamlit as st


# ─── Azure env check ─────────────────────────────────────────────────────────

def azure_env_ok() -> bool:
    return all(
        [
            os.getenv("AZURE_OPENAI_ENDPOINT"),
            os.getenv("AZURE_OPENAI_KEY"),
            os.getenv("AZURE_OPENAI_DEPLOYMENT"),
            os.getenv("AZURE_OPENAI_API_VERSION"),
        ]
    )


# ─── LangGraph agent ─────────────────────────────────────────────────────────

try:
    from typing import TypedDict
except ImportError:
    from typing_extensions import TypedDict


class AgentState(TypedDict):
    df: pd.DataFrame
    hierarchy_cols: List[str]
    has_variance_col: bool
    variance_col: str
    base_scenario: str
    compare_scenario: str
    path_trace: List[str]
    final_level_data: List[str]
    tree_data: List[Dict[str, Any]]
    final_summary: str


def calculate_variance_node(state: AgentState) -> Dict[str, Any]:
    df        = state["df"].copy()
    hierarchy = state.get("hierarchy_cols", [])
    has_var   = state.get("has_variance_col", True)

    if not hierarchy:
        return {"path_trace": ["Error: No hierarchy columns selected."], "final_level_data": [], "tree_data": []}

    if has_var:
        target_col = state.get("variance_col")
        if target_col not in df.columns:
            return {"path_trace": [f"Error: Column '{target_col}' not found."], "final_level_data": [], "tree_data": []}
    else:
        base_col = state.get("base_scenario")
        comp_col = state.get("compare_scenario")
        if base_col not in df.columns or comp_col not in df.columns:
            return {"path_trace": ["Error: Scenario columns not found."], "final_level_data": [], "tree_data": []}
        target_col = "Calculated_Variance"
        df[target_col] = (
            pd.to_numeric(df[base_col], errors="coerce").fillna(0)
            - pd.to_numeric(df[comp_col], errors="coerce").fillna(0)
        )

    df[target_col] = pd.to_numeric(df[target_col], errors="coerce").fillna(0)

    def format_variance(value: float) -> str:
        return f"{value / 1e6:,.2f}M"

    total_variance = df[target_col].fillna(0).sum()
    path_trace       = [f"Overall Total Variance: {format_variance(total_variance)}"]
    final_level_data = [f"Overall Total Variance: {format_variance(total_variance)}"]

    def build_tree(current_df, depth) -> Tuple[List, List, List]:
        if depth >= len(hierarchy) or current_df.empty:
            return [], [], []
        col_name = hierarchy[depth]
        is_first = depth == 0
        is_last  = depth == len(hierarchy) - 1
        grouped  = current_df.groupby(col_name)[target_col].sum()
        if grouped.empty or grouped.isna().all():
            return [], [], []
        top_5 = grouped.reindex(grouped.abs().sort_values(ascending=False).index).head(5)
        trace_lines, final_lines, tree_nodes = [], [], []

        for item, val in top_5.items():
            if pd.isna(item):
                continue
            item_label  = str(item)
            value_label = format_variance(float(val))

            if is_first:
                trace_lines.append(f"Primary Category: '{item_label}' (Total: {value_label})")
                final_lines.append(f"\nPrimary Category: '{item_label}' (Total Variance: {value_label})")
                title = f"Primary Category: {item_label} ({value_label})"
            elif is_last:
                trace_lines.append(f"Final Level ({col_name}): '{item_label}' -> {value_label}")
                final_lines.append(f"  - {col_name} '{item_label}': {value_label}")
                title = f"Final Level | {col_name}: {item_label} ({value_label})"
            else:
                trace_lines.append(f"Driver ({col_name}): '{item_label}' -> {value_label}")
                title = f"Driver | {col_name}: {item_label} ({value_label})"

            node: Dict[str, Any] = {
                "column": col_name, "item": item_label,
                "value": float(val), "value_display": value_label,
                "title": title, "children": [],
            }
            if not is_last:
                next_df = current_df[current_df[col_name] == item]
                sub_trace, sub_final, sub_nodes = build_tree(next_df, depth + 1)
                trace_lines.extend(sub_trace)
                final_lines.extend(sub_final)
                node["children"] = sub_nodes
            tree_nodes.append(node)
        return trace_lines, final_lines, tree_nodes

    tree_trace, tree_final, tree_nodes = build_tree(df, 0)
    path_trace.extend(tree_trace)
    final_level_data.extend(tree_final)
    return {"path_trace": path_trace, "final_level_data": final_level_data, "tree_data": tree_nodes}


def synthesize_insight_node(state: AgentState) -> Dict[str, Any]:
    if state["path_trace"] and "Error:" in state["path_trace"][0]:
        return {"final_summary": "Analysis aborted due to invalid data configuration."}

    from langchain_core.messages import HumanMessage, SystemMessage
    from langchain_openai import AzureChatOpenAI

    llm = AzureChatOpenAI(
        azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
        api_key=os.getenv("AZURE_OPENAI_KEY"),
        azure_deployment=os.getenv("AZURE_OPENAI_DEPLOYMENT"),
        api_version=os.getenv("AZURE_OPENAI_API_VERSION"),
    )

    system_prompt = (
        "You are a strict, professional financial data analyst. Provide an Executive Summary formatted EXACTLY as follows:\n"
        "1. A brief 1-2 sentence overall conclusion regarding the total variance.\n"
        "2. A bulleted breakdown for each 'Primary Category' analyzed. Under each, list the Top 5 reasons/drivers provided, "
        "along with exact variance amounts.\n\n"
        "---ROOT CAUSE ANALYSIS---\n"
        "Provide a detailed 4 to 5 line analytical paragraph explaining the underlying root causes of the variance "
        "based strictly on the provided data.\n\n"
        "---CATEGORY COMMENTARY---\n"
        "Act as a Senior Business Partner. For each 'Primary Category', DO NOT just repeat the numbers. "
        "Deduce the operational and financial reasons behind the variance based on the sub-categories provided. "
        "Provide 4 to 5 deep insights formatted EXACTLY as:\n"
        "[Primary Category Name]:\n"
        "- Operational Driver: [Specific business activity or market condition]\n"
        "- Financial Impact: [How this affects profit margins or budgets]\n"
        "- Risk/Opportunity: [Growing risk or area to capitalize on]\n"
        "- Recommended Action: [Specific next step or mitigation strategy]\n"
        "Do not add conversational filler."
    )

    messages = [
        SystemMessage(content=system_prompt),
        HumanMessage(content=f"Filtered Final Level Data:\n{chr(10).join(state['final_level_data'])}"),
    ]
    response = llm.invoke(messages)
    return {"final_summary": response.content}


def build_graph():
    from langgraph.graph import END, START, StateGraph

    workflow = StateGraph(AgentState)
    workflow.add_node("calculate_variance", calculate_variance_node)
    workflow.add_node("synthesize_insight", synthesize_insight_node)
    workflow.add_edge(START, "calculate_variance")
    workflow.add_edge("calculate_variance", "synthesize_insight")
    workflow.add_edge("synthesize_insight", END)
    return workflow.compile()


def count_leaf_nodes(nodes: List[Dict[str, Any]]) -> int:
    count = 0
    for node in nodes:
        children = node.get("children", [])
        if children:
            count += count_leaf_nodes(children)
        else:
            count += 1
    return count


def render_trace_tree(nodes: List[Dict[str, Any]]) -> None:
    for node in nodes:
        if node.get("children"):
            with st.expander(node["title"], expanded=False):
                render_trace_tree(node["children"])
        else:
            st.markdown(f"- **{node['title']}**")


# ─── PPTX generator ──────────────────────────────────────────────────────────

def clean_markdown_for_ppt(text: str) -> str:
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*",     r"\1", text)
    text = re.sub(r"#(.*?)\n",      r"\1\n", text)
    return text.strip()


def generate_ppt_deck(
    total_variance: str,
    exec_summary: str,
    rca_text: str,
    comm_text: str,
    tree_data: List[Dict[str, Any]],
) -> io.BytesIO:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    primary_color = RGBColor(14, 43, 92)
    accent_color  = RGBColor(0, 163, 224)
    text_color    = RGBColor(89, 89, 89)
    blank_layout  = prs.slide_layouts[6]

    def add_header(slide, title_text):
        header = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        hp = header.text_frame.paragraphs[0]
        hp.text = title_text
        hp.font.size = Pt(28)
        hp.font.bold = True
        hp.font.color.rgb = primary_color
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(1.1), Inches(9.5), Inches(1.1)
        )
        line.line.color.rgb = accent_color
        line.line.width = Pt(2)

    # Slide 1 — Title
    s1 = prs.slides.add_slide(blank_layout)
    banner = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
    banner.fill.solid()
    banner.fill.fore_color.rgb = primary_color
    banner.line.color.rgb = primary_color
    tf = s1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2)).text_frame
    p = tf.paragraphs[0]
    p.text = "Variance Analysis & Root Cause Report"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = f"Total Impact: {total_variance}\nReport Date: {datetime.now().strftime('%B %d, %Y')}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = text_color
    p2.alignment = PP_ALIGN.CENTER

    # Slide 2 — Executive Summary
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "Executive Summary")
    cf = s2.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5)).text_frame
    cf.word_wrap = True
    clean_sum = clean_markdown_for_ppt(exec_summary)
    lines = [l.strip() for l in clean_sum.split("\n") if l.strip()]
    for idx, para in enumerate(lines):
        p = cf.paragraphs[0] if idx == 0 else cf.add_paragraph()
        is_bullet = para.startswith("-") or para.startswith("*")
        clean_text = para.lstrip("-*1234567890. ")
        if is_bullet:
            p.text = f"  •  {clean_text}"
            p.font.size = Pt(14)
            p.font.bold = False
            p.font.color.rgb = text_color
        else:
            p.text = clean_text
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = primary_color

    # Slide 3 — Drill-Down Drivers
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Recursive Drill-Down Drivers")
    df_frame = s3.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5)).text_frame
    df_frame.word_wrap = True
    for idx, node in enumerate(tree_data):
        p = df_frame.paragraphs[0] if idx == 0 else df_frame.add_paragraph()
        p.text = f"{node['item']} (Impact: {node['value_display']})"
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = primary_color
        for child in node.get("children", [])[:4]:
            p_sub = df_frame.add_paragraph()
            p_sub.text = f"  •  {child['column']} - {child['item']} ({child['value_display']})"
            p_sub.font.size = Pt(14)
            p_sub.font.bold = False
            p_sub.font.color.rgb = text_color

    # Slide 4 — Root Cause Analysis
    if rca_text and "generation failed" not in rca_text.lower():
        s4 = prs.slides.add_slide(blank_layout)
        add_header(s4, "Root Cause Analysis")
        rca_frame = s4.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5)).text_frame
        rca_frame.word_wrap = True
        p = rca_frame.paragraphs[0]
        p.text = clean_markdown_for_ppt(rca_text)
        p.font.size = Pt(16)
        p.font.bold = False
        p.font.color.rgb = text_color

    # Slide 5 — Category Commentary
    if comm_text and "generation failed" not in comm_text.lower():
        s5 = prs.slides.add_slide(blank_layout)
        add_header(s5, "Category Commentary")
        comm_frame = s5.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5)).text_frame
        comm_frame.word_wrap = True
        comm_lines = [l.strip() for l in clean_markdown_for_ppt(comm_text).split("\n") if l.strip()]
        for idx, para in enumerate(comm_lines):
            p = comm_frame.paragraphs[0] if idx == 0 else comm_frame.add_paragraph()
            if para.startswith("-") or para.startswith("*"):
                p.text = f"  •  {para.lstrip('-* ')}"
                p.font.size = Pt(14)
                p.font.bold = False
                p.font.color.rgb = text_color
            else:
                p.text = para
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = primary_color

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream
