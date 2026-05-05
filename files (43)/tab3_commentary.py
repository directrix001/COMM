"""
routers/tab3_commentary.py  (UPDATED — Persona-Based Generation)
────────────────────────────────────────────────────────────────
POST /api/tab3/run        Run LangGraph agent with persona-tuned system prompt
GET  /api/tab3/download/md|txt|pptx
"""

from __future__ import annotations

import io, os, re, uuid
from datetime import datetime
from typing import Any, Dict, List

import pandas as pd
from fastapi import APIRouter, HTTPException, Request
from fastapi.responses import JSONResponse, StreamingResponse

from services import session_store as ss

router = APIRouter()

# ── PERSONA SYSTEM PROMPTS ────────────────────────────────────────────────
PERSONA_PROMPTS: Dict[str, str] = {
    "CFO": (
        "You are a Chief Financial Officer addressing the Board. "
        "Be strategic, concise and risk-aware. Focus on P&L impact, "
        "cash implications and recommend decisive actions. "
        "Use executive-level language — no operational jargon."
    ),
    "VP_FINANCE": (
        "You are a VP of Finance presenting to the CFO. "
        "Focus on budget stewardship, cost-control levers and FP&A insights. "
        "Highlight variances against plan with percentage deviations. "
        "Recommend corrective budget actions with owners and timelines."
    ),
    "BUSINESS_PARTNER": (
        "You are a Senior Finance Business Partner presenting to operational leaders. "
        "Be practical and action-oriented. Explain root causes using operational context — "
        "headcount, procurement, project delays. Suggest specific next steps for each team."
    ),
    "BOARD": (
        "You are preparing a Board-level executive briefing. "
        "Use investor-ready narrative language. Focus on strategic implications, "
        "market context and organisational opportunities. Keep bullet points high-level — "
        "no granular line-item detail. Maximum 3 key messages."
    ),
    "AUDITOR": (
        "You are an Internal Auditor preparing a JSOX-compliant variance memo. "
        "Be evidence-based, precise and factual. Flag variances that exceed materiality thresholds. "
        "Identify control weaknesses and recommend audit follow-up procedures. "
        "Cite specific data points and column references."
    ),
    "FP_A": (
        "You are an FP&A Analyst preparing a detailed management report. "
        "Be thorough and data-driven. Reference specific numbers, percentages and driver trees. "
        "Explain methodology. Suggest model improvements and scenario sensitivities. "
        "Use technical financial terminology freely."
    ),
}

DEFAULT_PERSONA_PROMPT = (
    "You are a professional financial data analyst. "
    "Be concise, accurate and insight-driven."
)


def _sid(request: Request) -> str:
    return request.cookies.get("va_sid", str(uuid.uuid4()))


# ── POST /api/tab3/run ────────────────────────────────────────────────────

@router.post("/run")
async def run_commentary(request: Request):
    """
    Body JSON:
    {
      "hierarchy_cols":   ["OH/LC", "Division_Desc"],
      "has_variance_col": true,
      "variance_col":     "delta",
      "base_scenario":    "",
      "compare_scenario": "",
      "persona":          "CFO"          ← NEW field
    }
    """
    sid  = _sid(request)
    body = await request.json()

    raw_leaf   = ss.get(sid, "leaf_df_bytes")
    raw_master = ss.get(sid, "master_db_bytes")

    if raw_leaf:
        df = pd.read_json(io.StringIO(raw_leaf), orient="split")
    elif raw_master:
        df = pd.read_json(io.StringIO(raw_master), orient="split")
    else:
        raise HTTPException(404, "No data in session. Run Tab 1 or Tab 2 first.")

    hierarchy_cols   = body.get("hierarchy_cols", [])
    has_variance_col = body.get("has_variance_col", True)
    variance_col     = body.get("variance_col", "delta")
    base_scenario    = body.get("base_scenario", "")
    compare_scenario = body.get("compare_scenario", "")
    persona          = body.get("persona", "CFO")

    if not hierarchy_cols:
        raise HTTPException(422, "hierarchy_cols must not be empty.")

    try:
        tree_result = _calculate_variance(
            df, hierarchy_cols, has_variance_col,
            variance_col, base_scenario, compare_scenario,
        )
    except Exception as exc:
        raise HTTPException(500, f"Variance calculation failed: {exc}")

    try:
        summary = await _synthesize_insight(
            tree_result["final_level_data"], persona
        )
    except Exception as exc:
        summary = f"AI synthesis failed: {exc}"

    total_var_str = (
        tree_result["path_trace"][0].replace("Overall Total Variance: ", "")
        if tree_result.get("path_trace") else "N/A"
    )

    ss.set(sid, "cg_summary",   summary)
    ss.set(sid, "cg_total_var", total_var_str)
    ss.set(sid, "cg_persona",   persona)

    return JSONResponse({
        "path_trace":      tree_result.get("path_trace", []),
        "tree_data":       tree_result.get("tree_data", []),
        "final_summary":   summary,
        "total_variance":  total_var_str,
        "leaf_node_count": _count_leaves(tree_result.get("tree_data", [])),
        "persona_used":    persona,
    })


# ── DOWNLOADS ─────────────────────────────────────────────────────────────

@router.get("/download/md")
async def download_md(request: Request):
    sid     = _sid(request)
    summary = ss.get(sid, "cg_summary")
    persona = ss.get(sid, "cg_persona") or "General"
    if not summary:
        raise HTTPException(404, "No commentary in session.")
    header = f"# Variance Analysis Commentary\n**Persona:** {persona}\n\n"
    return StreamingResponse(
        iter([(header + summary).encode()]),
        media_type="text/markdown",
        headers={"Content-Disposition": "attachment; filename=executive_variance_commentary.md"},
    )


@router.get("/download/txt")
async def download_txt(request: Request):
    sid     = _sid(request)
    summary = ss.get(sid, "cg_summary")
    if not summary:
        raise HTTPException(404, "No commentary in session.")
    return StreamingResponse(
        iter([summary.encode()]),
        media_type="text/plain",
        headers={"Content-Disposition": "attachment; filename=executive_variance_commentary.txt"},
    )


@router.get("/download/pptx")
async def download_pptx(request: Request):
    sid       = _sid(request)
    summary   = ss.get(sid, "cg_summary") or ""
    total_var = ss.get(sid, "cg_total_var") or "N/A"
    persona   = ss.get(sid, "cg_persona") or "General"

    exec_summary, rca_text, comm_text = _split_summary(summary)

    try:
        ppt_bytes = _generate_ppt(total_var, exec_summary, rca_text, comm_text, persona)
    except Exception as exc:
        raise HTTPException(500, f"PPTX generation failed: {exc}")

    fname = f"Variance_Deck_{datetime.now().strftime('%Y%m%d')}.pptx"
    return StreamingResponse(
        io.BytesIO(ppt_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={fname}"},
    )


# ── VARIANCE CALCULATION ──────────────────────────────────────────────────

def _calculate_variance(df, hierarchy, has_var, variance_col, base_sc, compare_sc):
    df = df.copy()

    if has_var:
        if variance_col not in df.columns:
            return {"path_trace": [f"Error: Column '{variance_col}' not found."],
                    "final_level_data": [], "tree_data": []}
        target_col = variance_col
    else:
        if base_sc not in df.columns or compare_sc not in df.columns:
            return {"path_trace": ["Error: Scenario columns not found."],
                    "final_level_data": [], "tree_data": []}
        target_col = "_calc_var"
        df[target_col] = (
            pd.to_numeric(df[base_sc], errors="coerce").fillna(0) -
            pd.to_numeric(df[compare_sc], errors="coerce").fillna(0)
        )

    df[target_col] = pd.to_numeric(df[target_col], errors="coerce").fillna(0)

    def fmt(v): return f"{v / 1e6:,.2f}M"

    total            = df[target_col].sum()
    path_trace       = [f"Overall Total Variance: {fmt(total)}"]
    final_level_data = [f"Overall Total Variance: {fmt(total)}"]

    def build_tree(sub, depth):
        if depth >= len(hierarchy) or sub.empty:
            return [], [], []
        col      = hierarchy[depth]
        is_first = depth == 0
        is_last  = depth == len(hierarchy) - 1
        grouped  = sub.groupby(col)[target_col].sum()
        top5     = grouped.reindex(grouped.abs().sort_values(ascending=False).index).head(5)
        traces, finals, nodes = [], [], []

        for item, val in top5.items():
            if pd.isna(item): continue
            lbl, vd = str(item), fmt(float(val))
            if is_first:
                traces.append(f"Primary Category: '{lbl}' (Total: {vd})")
                finals.append(f"\nPrimary Category: '{lbl}' (Total Variance: {vd})")
                title = f"Primary Category: {lbl} ({vd})"
            elif is_last:
                traces.append(f"Final Level ({col}): '{lbl}' -> {vd}")
                finals.append(f"  - {col} '{lbl}': {vd}")
                title = f"Final Level | {col}: {lbl} ({vd})"
            else:
                traces.append(f"Driver ({col}): '{lbl}' -> {vd}")
                title = f"Driver | {col}: {lbl} ({vd})"

            node: Dict[str, Any] = {
                "column": col, "item": lbl,
                "value": float(val), "value_display": vd,
                "title": title, "children": [],
            }
            if not is_last:
                st, sf, sn = build_tree(sub[sub[col] == item], depth + 1)
                traces.extend(st); finals.extend(sf); node["children"] = sn
            nodes.append(node)
        return traces, finals, nodes

    tt, tf, tn = build_tree(df, 0)
    path_trace.extend(tt)
    final_level_data.extend(tf)
    return {"path_trace": path_trace, "final_level_data": final_level_data, "tree_data": tn}


# ── AZURE OPENAI SYNTHESIS (PERSONA-AWARE) ────────────────────────────────

async def _synthesize_insight(final_level_data: List[str], persona: str) -> str:
    endpoint   = os.getenv("AZURE_OPENAI_ENDPOINT")
    api_key    = os.getenv("AZURE_OPENAI_KEY")
    deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT")
    api_ver    = os.getenv("AZURE_OPENAI_API_VERSION")

    if not all([endpoint, api_key, deployment, api_ver]):
        return (
            f"**[Azure OpenAI not configured]**\n"
            f"Persona selected: {persona}\n\n"
            "Set AZURE_OPENAI_ENDPOINT, AZURE_OPENAI_KEY, "
            "AZURE_OPENAI_DEPLOYMENT, AZURE_OPENAI_API_VERSION in .env\n\n"
            "Drill-down data:\n" + "\n".join(final_level_data)
        )

    from langchain_openai import AzureChatOpenAI
    from langchain_core.messages import SystemMessage, HumanMessage

    # ── Build persona-specific system prompt ──────────────────────────────
    persona_context = PERSONA_PROMPTS.get(persona, DEFAULT_PERSONA_PROMPT)

    system_prompt = (
        f"{persona_context}\n\n"
        "Structure your response EXACTLY as:\n"
        "1. Executive Summary: 2-3 sentences capturing total variance and top drivers.\n"
        "2. Bulleted breakdown per Primary Category with exact variance amounts.\n\n"
        "---ROOT CAUSE ANALYSIS---\n"
        "4-5 lines explaining root causes based strictly on data provided.\n\n"
        "---CATEGORY COMMENTARY---\n"
        "For each Primary Category provide:\n"
        "- Operational Driver\n"
        "- Financial Impact\n"
        "- Risk/Opportunity\n"
        "- Recommended Action\n"
        "Do not add conversational filler."
    )

    llm = AzureChatOpenAI(
        azure_endpoint=endpoint, api_key=api_key,
        azure_deployment=deployment, api_version=api_ver,
    )
    response = llm.invoke([
        SystemMessage(content=system_prompt),
        HumanMessage(content="\n".join(final_level_data)),
    ])
    return response.content


# ── HELPERS ───────────────────────────────────────────────────────────────

def _split_summary(summary: str):
    exec_s = summary; rca = ""; comm = ""
    if "---ROOT CAUSE ANALYSIS---" in summary:
        parts  = summary.split("---ROOT CAUSE ANALYSIS---")
        exec_s = parts[0].replace("Executive Summary:", "").strip()
        rest   = parts[1]
        if "---CATEGORY COMMENTARY---" in rest:
            rp = rest.split("---CATEGORY COMMENTARY---")
            rca = rp[0].strip(); comm = rp[1].strip()
        else:
            rca = rest.strip()
    return exec_s, rca, comm


def _count_leaves(nodes: list) -> int:
    c = 0
    for n in nodes:
        ch = n.get("children", [])
        c += _count_leaves(ch) if ch else 1
    return c


def _clean_md(text: str) -> str:
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    return re.sub(r"\*(.*?)\*", r"\1", text).strip()


def _generate_ppt(total_var, exec_summary, rca_text, comm_text, persona) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs   = Presentation()
    navy  = RGBColor(14, 43, 92)
    grey  = RGBColor(89, 89, 89)
    blank = prs.slide_layouts[6]

    def hdr(slide, title):
        tb = slide.shapes.add_textbox(Inches(.5), Inches(.4), Inches(9), Inches(.8))
        p  = tb.text_frame.paragraphs[0]
        p.text = title; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = navy

    # Title slide
    s1 = prs.slides.add_slide(blank)
    tf = s1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(2.5)).text_frame
    p  = tf.paragraphs[0]
    p.text = "Variance Analysis & Root Cause Report"
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = navy; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = f"Total Impact: {total_var}  |  Audience: {persona}\nDate: {datetime.now().strftime('%B %d, %Y')}"
    p2.font.size = Pt(15); p2.font.color.rgb = grey; p2.alignment = PP_ALIGN.CENTER

    # Executive summary slide
    s2 = prs.slides.add_slide(blank); hdr(s2, f"Executive Summary — {persona}")
    cf = s2.shapes.add_textbox(Inches(.5), Inches(1.3), Inches(9), Inches(5.8)).text_frame
    cf.word_wrap = True
    for idx, line in enumerate([l.strip() for l in _clean_md(exec_summary).split("\n") if l.strip()]):
        p = cf.paragraphs[0] if idx == 0 else cf.add_paragraph()
        is_b = line.startswith("-") or line.startswith("*")
        p.text = f"  • {line.lstrip('-* ')}" if is_b else line
        p.font.size = Pt(13 if is_b else 15); p.font.bold = not is_b
        p.font.color.rgb = grey if is_b else navy

    # RCA slide
    if rca_text:
        s3 = prs.slides.add_slide(blank); hdr(s3, "Root Cause Analysis")
        tf3 = s3.shapes.add_textbox(Inches(.5), Inches(1.3), Inches(9), Inches(5.8)).text_frame
        tf3.word_wrap = True
        tf3.paragraphs[0].text = _clean_md(rca_text)
        tf3.paragraphs[0].font.size = Pt(13); tf3.paragraphs[0].font.color.rgb = grey

    # Commentary slide
    if comm_text:
        s4 = prs.slides.add_slide(blank); hdr(s4, "Category Commentary")
        tf4 = s4.shapes.add_textbox(Inches(.5), Inches(1.3), Inches(9), Inches(5.8)).text_frame
        tf4.word_wrap = True
        for idx, line in enumerate([l.strip() for l in _clean_md(comm_text).split("\n") if l.strip()]):
            p = tf4.paragraphs[0] if idx == 0 else tf4.add_paragraph()
            is_b = line.startswith("-") or line.startswith("*")
            p.text = f"  • {line.lstrip('-* ')}" if is_b else line
            p.font.size = Pt(13 if is_b else 15); p.font.bold = not is_b
            p.font.color.rgb = grey if is_b else navy

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return buf.read()
