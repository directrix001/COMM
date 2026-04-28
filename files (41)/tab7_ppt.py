"""
routers/tab7_ppt.py
────────────────────
Tab 7 — PPT Upload & Master Database Push (FastAPI router)

Endpoints
---------
POST /api/tab7/upload       Upload PPT → extract → enrich → return preview
POST /api/tab7/push         Push current session data → master DB
GET  /api/tab7/download     Download enriched Excel
GET  /api/tab7/master       Preview master DB (first 100 rows)
"""

from __future__ import annotations

import io
import os
import re
import uuid
from datetime import datetime
from pathlib import Path

import pandas as pd
from fastapi import APIRouter, File, HTTPException, Request, UploadFile
from fastapi.responses import JSONResponse, StreamingResponse

from services import session_store as ss

router = APIRouter()

MASTER_XLSX = Path("database") / "Segregateddata.xlsx"


def _sid(request: Request) -> str:
    return request.cookies.get("va_sid", str(uuid.uuid4()))


# ── Reference lists ──────────────────────────────────────────────────────────

REGIONS = [
    'Afghanistan','Albania','Algeria','Andorra','Angola','Argentina','Armenia',
    'Australia','Austria','Azerbaijan','Bahrain','Bangladesh','Belgium','Brazil',
    'Bulgaria','Canada','Chile','China','Colombia','Croatia','Cuba','Cyprus',
    'Denmark','Egypt','Estonia','Ethiopia','Finland','France','Germany','Ghana',
    'Greece','Hungary','Iceland','India','Indonesia','Iran','Iraq','Ireland',
    'Israel','Italy','Jamaica','Japan','Jordan','Kazakhstan','Kenya','Kuwait',
    'Latvia','Lebanon','Lithuania','Luxembourg','Malaysia','Malta','Mexico',
    'Morocco','Netherlands','New Zealand','Nigeria','Norway','Oman','Pakistan',
    'Philippines','Poland','Portugal','Qatar','Romania','Russia','Saudi Arabia',
    'Serbia','Singapore','Slovakia','Slovenia','South Africa','South Korea',
    'Spain','Sri Lanka','Sweden','Switzerland','Syria','Thailand','Tunisia',
    'Turkey','Ukraine','United Arab Emirates','United Kingdom',
    'United States of America','Vietnam','Zimbabwe',
    'AMIO','AMIEO','Middle East','Europe','Oceania','ME','Africa',
]
LC_OH    = ['LC','OH','Notes','OVH']
CRITERIA = [
    'PROCURED SERVICES','TRAVEL & MEALS','EMPLOYEE WELFARE',
    'RECHARGE NISSAN Level0','OPERATING COSTS','OFFICE SPACE',
    'EMPLOYEE ACTIVITY COSTS','TAX','RECHARGE OUTSIDE',
    'PROVISION FOR DOUBTFUL DEBTS','COMPANY CAR COSTS','DEPRECIATION',
]
KEYWORDS = [
    'FINANCE & ACCOUNTING','TREASURY','TAX','SSC','LEGAL','COMPLIANCE',
    'COMMUNICATION','PURCHASING','AFTER SALES OPS','MARCOM',
    'CONNECTED SERVICES','MARKET INTELLIGENCE','ELECTRIFICATION','General Affairs',
]
ENTITY = [
    'Nissan Automotive Europe','NIBSA','NMISA','NMUK - PLANT',
    'NCE Germany','NITA','NMGB','NNE','Nissan France',
    'Nissan International SA','NAE','NRBS','NTCE','NMEF',
]


def _match(text: str, lst: list) -> str:
    if not isinstance(text, str): return ""
    tl = text.lower()
    return ", ".join(i for i in lst if re.search(r'\b' + re.escape(i.lower()) + r'\b', tl))


# ── POST /api/tab7/upload ────────────────────────────────────────────────────

@router.post("/upload")
async def upload_ppt(request: Request, file: UploadFile = File(...)):
    sid = _sid(request)
    if not file.filename.lower().endswith((".pptx", ".ppt")):
        raise HTTPException(400, "Only .pptx / .ppt files are accepted.")

    raw = await file.read()
    try:
        df = _extract_ppt(io.BytesIO(raw), file.filename)
    except Exception as exc:
        raise HTTPException(500, f"PPT extraction failed: {exc}")

    if df.empty:
        raise HTTPException(422, "No comments longer than 35 characters found in this file.")

    ss.set(sid, "ppt_extract_bytes", df.to_json(orient="split"))

    resp = JSONResponse({
        "session_id": sid,
        "rows":    len(df),
        "columns": df.columns.tolist(),
        "records": df.where(pd.notnull(df), None).to_dict(orient="records"),
    })
    resp.set_cookie("va_sid", sid, httponly=True, samesite="lax")
    return resp


# ── POST /api/tab7/push ──────────────────────────────────────────────────────

@router.post("/push")
async def push_to_master(request: Request):
    sid = _sid(request)
    raw = ss.get(sid, "ppt_extract_bytes")
    if not raw:
        raise HTTPException(404, "No extracted data in session. Upload a PPT first.")

    df_new = pd.read_json(io.StringIO(raw), orient="split")
    MASTER_XLSX.parent.mkdir(parents=True, exist_ok=True)

    if MASTER_XLSX.exists():
        master = pd.read_excel(MASTER_XLSX)
    else:
        master = pd.DataFrame()

    combined = pd.concat([master, df_new], ignore_index=True) if not master.empty else df_new.copy()
    if "Slide Number" in combined.columns and "Comments" in combined.columns:
        combined = combined.drop_duplicates(
            subset=["Slide Number","Comments"], keep="last"
        ).reset_index(drop=True)

    combined.to_excel(MASTER_XLSX, index=False)
    return JSONResponse({"status": "ok", "pushed_rows": len(df_new), "master_rows": len(combined)})


# ── GET /api/tab7/download ───────────────────────────────────────────────────

@router.get("/download")
async def download(request: Request):
    sid = _sid(request)
    raw = ss.get(sid, "ppt_extract_bytes")
    if not raw:
        raise HTTPException(404, "No extracted data in session.")
    df = pd.read_json(io.StringIO(raw), orient="split")
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as w:
        df.to_excel(w, index=False, sheet_name="Commentary_edited")
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=Commentary_edited.xlsx"},
    )


# ── GET /api/tab7/master ─────────────────────────────────────────────────────

@router.get("/master")
async def master_preview():
    if not MASTER_XLSX.exists():
        return JSONResponse({"columns": [], "records": [], "total": 0})
    df = pd.read_excel(MASTER_XLSX)
    preview = df.head(100)
    return JSONResponse({
        "columns": df.columns.tolist(),
        "records": preview.where(pd.notnull(preview), None).to_dict(orient="records"),
        "total":   len(df),
    })


# ── PPT extraction logic ─────────────────────────────────────────────────────

def _extract_ppt(file_obj, filename: str) -> pd.DataFrame:
    from pptx import Presentation
    prs     = Presentation(file_obj)
    content = []

    for slide_num, slide in enumerate(prs.slides, 1):
        header = None
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if header is None and text.startswith("G&A Evolution") and text not in ["AMIEO","M€ @BP25FX"]:
                    header = text
                content.append({"Slide Number": slide_num, "Comments": text})

        if header:
            header = header.replace("-","–")
            parts  = header.split(" – ")
            category = parts[0] if parts else ""
            month = scenario = fn = fn_view = ""

            m = re.search(
                r"\b(January|February|March|April|May|June|July|August|September"
                r"|October|November|December|Jan|Feb|Mar|Apr|May|Jun|Jul|Aug"
                r"|Sep|Oct|Nov|Dec)\b", header, re.IGNORECASE)
            if m: month = m.group(0)

            m2 = re.search(r"(MTD|YTD)\s+vs\.\s+[^–]+", header)
            if m2: scenario = m2.group(0)
            if len(parts) > 2: fn      = parts[2]
            if len(parts) > 3: fn_view = parts[3]

            for item in content:
                if item["Slide Number"] == slide_num:
                    item.update({
                        "File_name": filename, "Category": category,
                        "Scenarios": scenario, "Functions": fn,
                        "Functions-View": fn_view, "Month": month,
                        "Year": datetime.now().year,
                        "Forecast": "", "Actual": "", "Variance": "",
                    })

    filtered = [i for i in content if len(i.get("Comments","")) > 35]
    df = pd.DataFrame(filtered)
    if df.empty: return df

    df["CostCat description"] = df["Comments"].apply(lambda t: _match(t, CRITERIA))
    df["Function_desc"]       = df["Comments"].apply(lambda t: _match(t, KEYWORDS))
    df["Entity_desc"]         = df["Comments"].apply(lambda t: _match(t, ENTITY))
    df["Criteria"]            = df["Comments"].apply(lambda t: _match(t, LC_OH))
    df["Region"]              = df["Comments"].apply(lambda t: _match(t, REGIONS))
    return df
