"""
routers/tab3_commentary.py
───────────────────────────
Tab 3 — AI Commentary Generator (FastAPI router)

Endpoints
---------
POST /api/tab3/run       Run LangGraph agent → executive summary + RCA + category commentary
GET  /api/tab3/download/md    Download .md report
GET  /api/tab3/download/txt   Download .txt report
GET  /api/tab3/download/pptx  Download .pptx deck
"""

from __future__ import annotations

import io
import os
import re
import uuid
from datetime import datetime
from typing import Any, Dict, List

import pandas as pd
from fastapi import APIRouter, HTTPException, Request
from fastapi.responses import JSONResponse, StreamingResponse

from services import session_store as ss

router = APIRouter()


def _sid(request: Request) -> str:
    return request.cookies.get("va_sid", str(uuid.uuid4()))


# ── POST /api/tab3/run ───────────────────────────────────────────────────────

@router.post("/run")
async def run_commentary(request: Request):
    """
    Body JSON:
    {
      "hierarchy_cols":   ["OH/LC", "Division_Desc"],
      "has_variance_col": true,
      "variance_col":     "delta",
      "base_scenario":    "",
      "compare_scenario": ""
    }
    Uses leaf_df from session (set by Tab 2 run).
    """
    sid  = _sid(request)
    body = await request.json()

    # Load data — prefer Tab 2 leaf_df, fall back to master_db
    raw_leaf   = ss.get(sid, "leaf_df_bytes")
    raw_master = ss.get(sid, "master_db_bytes")

    if raw_leaf:
        df = pd.read_json(io.StringIO(raw_leaf), orient="split")
    elif raw_master:
        df = pd.read_json(io.StringIO(raw_master), orient="split")
    else:
        raise HTTPException(404, "No data in session. Run Tab 1 mapping or Tab 2 variance first.")

    hierarchy_cols   = body.get("hierarchy_cols", [])
    has_variance_col = body.get("has_variance_col", True)
    variance_col     = body.get("variance_col", "delta")
    base_scenario    = body.get("base_scenario", "")
    compare_scenario = body.get("compare_scenario", "")

    if not hierarchy_cols:
        raise HTTPException(422, "hierarchy_cols must not be empty.")

    # ── Calculate variance tree ──────────────────────────────────────────────
    try:
        tree_result = _calculate_variance(
            df, hierarchy_cols, has_variance_col,
            variance_col, base_scenario, compare_scenario
        )
    except Exception as exc:
        raise HTTPException(500, f"Variance calculation failed: {exc}")

    # ── Azure OpenAI synthesis ───────────────────────────────────────────────
    try:
        summary = await _synthesize_insight(tree_result["final_level_data"])
    except Exception as exc:
        summary = f"AI synthesis failed: {exc}"

    total_var_str = (
        tree_result["path_trace"][0].replace("Overall Total Variance: ", "")
        if tree_result.get("path_trace") else "N/A"
    )

    # Store for download
    ss.set(sid, "cg_summary",       summary)
    ss.set(sid, "cg_total_var",     total_var_str)
    ss.set(sid, "cg_tree_data",     str(tree_result.get("tree_data", [])))

    return JSONResponse({
        "path_trace":      tree_result.get("path_trace", []),
        "tree_data":       tree_result.get("tree_data", []),
        "final_summary":   summary,
        "total_variance":  total_var_str,
        "leaf_node_count": _count_leaves(tree_result.get("tree_data", [])),
    })


# ── GET downloads ────────────────────────────────────────────────────────────

@router.get("/download/md")
async def download_md(request: Request):
    sid     = _sid(request)
    summary = ss.get(sid, "cg_summary")
    if not summary:
        raise HTTPException(404, "No commentary in session.")
    return StreamingResponse(
        iter([summary.encode()]),
        media_type="text/markdown",
        headers={"Content-Disposition": "attachment; filename=executive_variance_commentary.md"},
    )


@router.get("/download/txt")
async def download_txt(request: Request):
    sid     = _sid(request)
    summary = ss.get(sid, "cg_summary")
    if not summary:
        raise HTTPException(404, "No commentary in session.")
    return StreamingResponse(
        iter([summary.encode()]),
        media_type="text/plain",
        headers={"Content-Disposition": "attachment; filename=executive_variance_commentary.txt"},
    )


@router.get("/download/pptx")
async def download_pptx(request: Request):
    sid       = _sid(request)
    summary   = ss.get(sid, "cg_summary") or ""
    total_var = ss.get(sid, "cg_total_var") or "N/A"

    exec_summary, rca_text, comm_text = _split_summary(summary)

    try:
        ppt_bytes = _generate_ppt(total_var, exec_summary, rca_text, comm_text, [])
    except Exception as exc:
        raise HTTPException(500, f"PPTX generation failed: {exc}")

    fname = f"Variance_Deck_{datetime.now().strftime('%Y%m%d')}.pptx"
    return StreamingResponse(
        io.BytesIO(ppt_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={fname}"},
    )


# ── Private helpers ──────────────────────────────────────────────────────────

def _calculate_variance(
    df: pd.DataFrame,
    hierarchy: List[str],
    has_var: bool,
    variance_col: str,
    base_scenario: str,
    compare_scenario: str,
) -> Dict[str, Any]:
    df = df.copy()

    if has_var:
        if variance_col not in df.columns:
            return {"path_trace": [f"Error: Column '{variance_col}' not found."],
                    "final_level_data": [], "tree_data": []}
        target_col = variance_col
    else:
        if base_scenario not in df.columns or compare_scenario not in df.columns:
            return {"path_trace": ["Error: Scenario columns not found."],
                    "final_level_data": [], "tree_data": []}
        target_col = "_calc_var"
        df[target_col] = (
            pd.to_numeric(df[base_scenario], errors="coerce").fillna(0) -
            pd.to_numeric(df[compare_scenario], errors="coerce").fillna(0)
        )

    df[target_col] = pd.to_numeric(df[target_col], errors="coerce").fillna(0)

    def fmt(v: float) -> str:
        return f"{v / 1e6:,.2f}M"

    total = df[target_col].sum()
    path_trace       = [f"Overall Total Variance: {fmt(total)}"]
    final_level_data = [f"Overall Total Variance: {fmt(total)}"]

    def build_tree(sub_df, depth):
        if depth >= len(hierarchy) or sub_df.empty:
            return [], [], []
        col      = hierarchy[depth]
        is_first = depth == 0
        is_last  = depth == len(hierarchy) - 1
        grouped  = sub_df.groupby(col)[target_col].sum()
        top5     = grouped.reindex(grouped.abs().sort_values(ascending=False).index).head(5)

        traces, finals, nodes = [], [], []
        for item, val in top5.items():
            if pd.isna(item):
                continue
            lbl = str(item); vd = fmt(float(val))
            if is_first:
                traces.append(f"Primary Category: '{lbl}' (Total: {vd})")
                finals.append(f"\nPrimary Category: '{lbl}' (Total Variance: {vd})")
                title = f"Primary Category: {lbl} ({vd})"
            elif is_last:
                traces.append(f"Final Level ({col}): '{lbl}' -> {vd}")
                finals.append(f"  - {col} '{lbl}': {vd}")
                title = f"Final Level | {col}: {lbl} ({vd})"
            else:
                traces.append(f"Driver ({col}): '{lbl}' -> {vd}")
                title = f"Driver | {col}: {lbl} ({vd})"

            node: Dict[str, Any] = {
                "column": col, "item": lbl,
                "value": float(val), "value_display": vd,
                "title": title, "children": [],
            }
            if not is_last:
                nxt = sub_df[sub_df[col] == item]
                st, sf, sn = build_tree(nxt, depth + 1)
                traces.extend(st); finals.extend(sf); node["children"] = sn
            nodes.append(node)
        return traces, finals, nodes

    tt, tf, tn = build_tree(df, 0)
    path_trace.extend(tt)
    final_level_data.extend(tf)
    return {"path_trace": path_trace, "final_level_data": final_level_data, "tree_data": tn}


async def _synthesize_insight(final_level_data: List[str]) -> str:
    """Call Azure OpenAI. Falls back to a formatted stub if env vars missing."""
    endpoint   = os.getenv("AZURE_OPENAI_ENDPOINT")
    api_key    = os.getenv("AZURE_OPENAI_KEY")
    deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT")
    api_ver    = os.getenv("AZURE_OPENAI_API_VERSION")

    if not all([endpoint, api_key, deployment, api_ver]):
        return (
            "**[Azure OpenAI not configured]**\n\n"
            "Set AZURE_OPENAI_ENDPOINT, AZURE_OPENAI_KEY, "
            "AZURE_OPENAI_DEPLOYMENT, AZURE_OPENAI_API_VERSION in your .env file.\n\n"
            "Drill-down data:\n" + "\n".join(final_level_data)
        )

    from langchain_openai import AzureChatOpenAI
    from langchain_core.messages import SystemMessage, HumanMessage

    llm = AzureChatOpenAI(
        azure_endpoint=endpoint, api_key=api_key,
        azure_deployment=deployment, api_version=api_ver,
    )
    system_prompt = (
        "You are a strict, professional financial data analyst. "
        "Provide an Executive Summary, then ---ROOT CAUSE ANALYSIS--- section, "
        "then ---CATEGORY COMMENTARY--- section. "
        "Be concise, data-driven, and avoid filler."
    )
    response = llm.invoke([
        SystemMessage(content=system_prompt),
        HumanMessage(content="\n".join(final_level_data)),
    ])
    return response.content


def _split_summary(summary: str):
    exec_s = summary; rca = ""; comm = ""
    if "---ROOT CAUSE ANALYSIS---" in summary:
        parts  = summary.split("---ROOT CAUSE ANALYSIS---")
        exec_s = parts[0].replace("Executive Summary:", "").strip()
        rest   = parts[1]
        if "---CATEGORY COMMENTARY---" in rest:
            rp = rest.split("---CATEGORY COMMENTARY---")
            rca = rp[0].strip(); comm = rp[1].strip()
        else:
            rca = rest.strip()
    return exec_s, rca, comm


def _count_leaves(nodes: list) -> int:
    c = 0
    for n in nodes:
        ch = n.get("children", [])
        c += _count_leaves(ch) if ch else 1
    return c


def _clean_md(text: str) -> str:
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*",     r"\1", text)
    return text.strip()


def _generate_ppt(total_var, exec_summary, rca_text, comm_text, tree_data) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs   = Presentation()
    navy  = RGBColor(14, 43, 92)
    cyan  = RGBColor(0, 163, 224)
    grey  = RGBColor(89, 89, 89)
    blank = prs.slide_layouts[6]

    def header(slide, title):
        tb = slide.shapes.add_textbox(Inches(.5), Inches(.4), Inches(9), Inches(1))
        p  = tb.text_frame.paragraphs[0]
        p.text = title; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = navy

    # Slide 1
    s1 = prs.slides.add_slide(blank)
    tf = s1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2)).text_frame
    p  = tf.paragraphs[0]
    p.text = "Variance Analysis & Root Cause Report"
    p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = navy; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = f"Total Impact: {total_var}\nReport Date: {datetime.now().strftime('%B %d, %Y')}"
    p2.font.size = Pt(16); p2.font.color.rgb = grey; p2.alignment = PP_ALIGN.CENTER

    # Slide 2 — exec summary
    s2 = prs.slides.add_slide(blank); header(s2, "Executive Summary")
    cf = s2.shapes.add_textbox(Inches(.5), Inches(1.4), Inches(9), Inches(5.5)).text_frame
    cf.word_wrap = True
    for idx, line in enumerate([l.strip() for l in _clean_md(exec_summary).split("\n") if l.strip()]):
        p = cf.paragraphs[0] if idx == 0 else cf.add_paragraph()
        is_b = line.startswith("-") or line.startswith("*")
        p.text = f"  • {line.lstrip('-* ')}" if is_b else line
        p.font.size = Pt(14 if is_b else 16)
        p.font.bold = not is_b
        p.font.color.rgb = grey if is_b else navy

    # Slide 3 — RCA
    if rca_text:
        s3 = prs.slides.add_slide(blank); header(s3, "Root Cause Analysis")
        tf3 = s3.shapes.add_textbox(Inches(.5), Inches(1.4), Inches(9), Inches(5.5)).text_frame
        tf3.word_wrap = True
        tf3.paragraphs[0].text = _clean_md(rca_text)
        tf3.paragraphs[0].font.size = Pt(14)
        tf3.paragraphs[0].font.color.rgb = grey

    buf = io.BytesIO()
    prs.save(buf); buf.seek(0)
    return buf.read()
